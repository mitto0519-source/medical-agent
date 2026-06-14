"""Universal file loader — 다양한 확장자를 텍스트(또는 vision-ready base64)로 변환.

오픈소스 기반 조합:
  - markitdown (Microsoft, MIT) — 90+ 형식 → markdown (best-effort 통합 진입)
  - pdfplumber (MIT) — PDF (markitdown fallback이 부족할 때)
  - python-docx / python-pptx / openpyxl (MIT) — Office 분리 핸들러
  - pyreadstat (Apache 2.0) — SPSS .sav / STATA .dta / SAS .sas7bdat
  - Pillow + base64 — 이미지(vision-ready data URI)
  - chardet — 인코딩 자동 감지 (raw 텍스트)

지원 확장자:
  텍스트       .txt .md .rst .log .csv .tsv .json .yaml .yml .xml .html .htm
  문서         .pdf .docx .doc .pptx .ppt .xlsx .xls .odt .epub .rtf
  데이터       .sav .dta .sas7bdat .parquet .feather .pkl
  코드         .py .js .ts .go .rs .java .cpp .c .sh .ipynb
  이미지       .png .jpg .jpeg .gif .webp .bmp .tiff (→ Vision LLM에 base64)
  음성·영상    .mp3 .wav .m4a .mp4 .mov (transcription은 미지원 — 첨부만 기록)
"""
from __future__ import annotations

import base64
import mimetypes
from pathlib import Path
from typing import Optional

from src.config.logging_config import get_logger

_log = get_logger(__name__)


_TEXT_EXTS = {".txt", ".md", ".rst", ".log", ".csv", ".tsv",
                ".json", ".yaml", ".yml", ".xml", ".html", ".htm",
                ".py", ".js", ".ts", ".go", ".rs", ".java", ".cpp",
                ".c", ".sh", ".bat", ".ps1", ".rb", ".php"}
_OFFICE_EXTS = {".pdf", ".docx", ".doc", ".pptx", ".ppt",
                  ".xlsx", ".xls", ".odt", ".odp", ".ods", ".epub", ".rtf"}
_DATA_EXTS = {".sav", ".dta", ".sas7bdat", ".parquet", ".feather", ".pkl", ".pickle"}
_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".tif"}
_MEDIA_EXTS = {".mp3", ".wav", ".m4a", ".ogg", ".flac",
                  ".mp4", ".mov", ".avi", ".mkv", ".webm"}
_NOTEBOOK_EXTS = {".ipynb"}


def detect_kind(path: Path) -> str:
    """Return one of: text / office / data / image / media / notebook / unknown."""
    ext = path.suffix.lower()
    if ext in _TEXT_EXTS:     return "text"
    if ext in _OFFICE_EXTS:   return "office"
    if ext in _DATA_EXTS:     return "data"
    if ext in _IMAGE_EXTS:    return "image"
    if ext in _MEDIA_EXTS:    return "media"
    if ext in _NOTEBOOK_EXTS: return "notebook"
    return "unknown"


def _read_text_raw(path: Path, max_bytes: int = 5_000_000) -> str:
    """Read text with chardet auto-encoding detection."""
    raw = path.read_bytes()[:max_bytes]
    try:
        import chardet
        enc = (chardet.detect(raw) or {}).get("encoding") or "utf-8"
    except Exception:
        enc = "utf-8"
    try:
        return raw.decode(enc, errors="replace")
    except Exception:
        return raw.decode("utf-8", errors="replace")


def _read_with_markitdown(path: Path) -> Optional[str]:
    """Microsoft markitdown — 단일 진입으로 90+ 형식 → markdown."""
    try:
        from markitdown import MarkItDown
        md = MarkItDown()
        result = md.convert(str(path))
        return result.text_content if hasattr(result, "text_content") else str(result)
    except ImportError:
        _log.debug("markitdown not installed; falling back to dedicated handlers")
        return None
    except Exception as e:
        _log.debug("markitdown convert fail (%s): %s", path.name, e)
        return None


def _read_pdf(path: Path) -> str:
    """PDF fallback chain — pdfplumber → pypdf → pdfminer.six."""
    for backend in ("pdfplumber", "pypdf", "pdfminer"):
        try:
            if backend == "pdfplumber":
                import pdfplumber
                with pdfplumber.open(str(path)) as pdf:
                    return "\n\n".join(p.extract_text() or "" for p in pdf.pages)
            elif backend == "pypdf":
                from pypdf import PdfReader
                rd = PdfReader(str(path))
                return "\n\n".join(p.extract_text() or "" for p in rd.pages)
            elif backend == "pdfminer":
                from pdfminer.high_level import extract_text
                return extract_text(str(path))
        except ImportError:
            continue
        except Exception as e:
            _log.debug("PDF %s backend fail: %s", backend, e)
            continue
    return ""


def _read_docx(path: Path) -> str:
    try:
        from docx import Document
        d = Document(str(path))
        return "\n\n".join(p.text for p in d.paragraphs if p.text.strip())
    except Exception as e:
        _log.debug("docx read fail: %s", e)
        return ""


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        p = Presentation(str(path))
        chunks = []
        for i, slide in enumerate(p.slides, 1):
            chunks.append(f"### Slide {i}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    chunks.append(shape.text)
        return "\n".join(chunks)
    except Exception as e:
        _log.debug("pptx read fail: %s", e)
        return ""


def _read_xlsx(path: Path) -> str:
    try:
        import pandas as pd
        sheets = pd.read_excel(str(path), sheet_name=None, nrows=200)
        out = []
        for name, df in sheets.items():
            out.append(f"### Sheet: {name} ({len(df)} rows × {len(df.columns)} cols)")
            out.append(df.head(30).to_markdown(index=False))
        return "\n\n".join(out)
    except Exception as e:
        _log.debug("xlsx read fail: %s", e)
        return ""


def _read_spss_sav(path: Path) -> str:
    try:
        import pyreadstat
        df, meta = pyreadstat.read_sav(str(path), metadataonly=False, row_limit=200)
        out = [f"## SPSS .sav — {len(df)} rows × {len(df.columns)} cols"]
        labels = (meta.column_names_to_labels or {})
        out.append("\n### Variable codebook (first 50)")
        for col in list(df.columns)[:50]:
            lbl = labels.get(col, "")
            out.append(f"- `{col}`: {lbl}")
        out.append("\n### Sample (first 20 rows)")
        out.append(df.head(20).to_markdown(index=False))
        return "\n".join(out)
    except Exception as e:
        _log.debug("sav read fail: %s", e)
        return ""


def _read_stata_dta(path: Path) -> str:
    try:
        import pyreadstat
        df, meta = pyreadstat.read_dta(str(path), metadataonly=False, row_limit=200)
        return f"## STATA .dta — {len(df)} rows × {len(df.columns)} cols\n\n" + \
               df.head(20).to_markdown(index=False)
    except Exception as e:
        try:
            import pandas as pd
            df = pd.read_stata(str(path), iterator=False, chunksize=None,
                                  convert_categoricals=False)
            return f"## STATA .dta — {len(df)} rows × {len(df.columns)} cols\n\n" + \
                   df.head(20).to_markdown(index=False)
        except Exception as e2:
            _log.debug("dta read fail: %s / %s", e, e2)
            return ""


def _read_csv(path: Path) -> str:
    try:
        import pandas as pd
        df = pd.read_csv(str(path), nrows=200)
        return (f"## CSV — {len(df)} rows × {len(df.columns)} cols (preview)\n\n" +
                df.head(30).to_markdown(index=False))
    except Exception as e:
        _log.debug("csv read fail: %s", e)
        return _read_text_raw(path)


def _read_notebook(path: Path) -> str:
    try:
        import json
        nb = json.loads(path.read_text(encoding="utf-8"))
        out = []
        for i, cell in enumerate(nb.get("cells", []), 1):
            src = "".join(cell.get("source", []))
            ctype = cell.get("cell_type", "?")
            out.append(f"### Cell {i} [{ctype}]\n```\n{src}\n```")
        return "\n\n".join(out)
    except Exception as e:
        _log.debug("ipynb read fail: %s", e)
        return _read_text_raw(path)


def _image_to_data_uri(path: Path, *, max_bytes: int = 5_000_000) -> str:
    """이미지를 base64 data URI로 — Vision LLM에 inline 전달용."""
    mime, _ = mimetypes.guess_type(str(path))
    if not mime:
        mime = f"image/{path.suffix.lstrip('.').lower() or 'png'}"
    data = path.read_bytes()[:max_bytes]
    return f"data:{mime};base64,{base64.b64encode(data).decode('ascii')}"


# ── Public API ──────────────────────────────────────────────────────────────

def load(path: str | Path) -> dict:
    """파일 → {kind, text, image_data_uri?, meta}.

    text는 LLM 컨텍스트에 직접 주입 가능한 markdown/plaintext.
    image_data_uri는 Vision LLM에 inline 전달 (Anthropic image content block).
    실패해도 dict는 항상 반환 (text="", error=...).
    """
    p = Path(path)
    if not p.exists() or not p.is_file():
        return {"path": str(p), "kind": "missing", "text": "",
                 "error": "file not found"}

    kind = detect_kind(p)
    meta = {"path": str(p), "name": p.name,
             "size_bytes": p.stat().st_size, "kind": kind}

    # 1) Images → vision data URI (text 추출 없음)
    if kind == "image":
        try:
            return {**meta, "text": f"[Image: {p.name} ({p.stat().st_size//1024}KB)]",
                    "image_data_uri": _image_to_data_uri(p)}
        except Exception as e:
            return {**meta, "text": "", "error": str(e)[:200]}

    # 2) Media (audio/video) — transcription 미지원, 메타만
    if kind == "media":
        return {**meta, "text": f"[Media file: {p.name} — transcription 미지원]"}

    # 3) Data files (SPSS/STATA/SAS/parquet) — dedicated handlers (markitdown 약함)
    if kind == "data":
        ext = p.suffix.lower()
        if ext == ".sav":  return {**meta, "text": _read_spss_sav(p)}
        if ext == ".dta":  return {**meta, "text": _read_stata_dta(p)}
        try:
            import pandas as pd
            if ext == ".parquet": df = pd.read_parquet(p)
            elif ext == ".feather": df = pd.read_feather(p)
            else: df = pd.read_pickle(p)
            return {**meta, "text": f"## {ext} — {len(df)} rows × {len(df.columns)} cols\n\n" +
                   df.head(20).to_markdown(index=False)}
        except Exception as e:
            return {**meta, "text": "", "error": str(e)[:200]}

    # 4) Notebooks
    if kind == "notebook":
        return {**meta, "text": _read_notebook(p)}

    # 5) Try markitdown first (broadest)
    text = _read_with_markitdown(p)
    if text:
        return {**meta, "text": text}

    # 6) Dedicated fallbacks per ext
    ext = p.suffix.lower()
    if ext == ".pdf":       return {**meta, "text": _read_pdf(p)}
    if ext == ".docx":      return {**meta, "text": _read_docx(p)}
    if ext == ".pptx":      return {**meta, "text": _read_pptx(p)}
    if ext in (".xlsx", ".xls"): return {**meta, "text": _read_xlsx(p)}
    if ext in (".csv", ".tsv"):  return {**meta, "text": _read_csv(p)}

    # 7) Plain text fallback (chardet auto-encoding)
    return {**meta, "text": _read_text_raw(p)}


def load_many(paths: list) -> list[dict]:
    return [load(p) for p in paths]


def render_for_llm(loaded: list[dict], *, max_text_per_file: int = 8000) -> str:
    """첨부 묶음을 system_prompt에 박을 수 있는 단일 블록으로."""
    if not loaded:
        return ""
    blocks = ["# 📎 ATTACHED FILES (사용자 첨부)"]
    for i, item in enumerate(loaded, 1):
        name = item.get("name", "?")
        kind = item.get("kind", "?")
        text = (item.get("text") or "")[:max_text_per_file]
        if not text:
            blocks.append(f"\n## [{i}] {name} ({kind}) — 본문 추출 실패")
            continue
        blocks.append(f"\n## [{i}] {name} ({kind})\n{text}")
        if len(item.get("text", "")) > max_text_per_file:
            blocks.append(f"... (잘림, 원본 {len(item['text']):,}자)")
    return "\n".join(blocks)


__all__ = ["load", "load_many", "render_for_llm",
            "detect_kind"]
