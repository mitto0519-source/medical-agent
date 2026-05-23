"""Universal document reader — PDF, Word, PowerPoint, Excel, text, images.

Supported formats:
  PDF     .pdf            → PyMuPDF
  Word    .docx           → python-docx
  PPT     .pptx           → python-pptx
  Excel   .xlsx .xls .csv → openpyxl / pandas
  Text    .txt .md .rst   → plain read
  Image   .jpg .jpeg .png .webp .bmp .gif → Claude Vision (중앙 모델 설정)
"""
from __future__ import annotations

import base64
import os
from pathlib import Path
from typing import Dict, List, Optional

from src.config.logging_config import get_logger
from src.config.models import get_vision_model

_log = get_logger(__name__)


# ── format readers ────────────────────────────────────────────────────────────

def _read_pdf(path: Path) -> tuple[str, dict]:
    import fitz
    doc = fitz.open(str(path))
    meta = doc.metadata or {}
    pages = [page.get_text("text").strip() for page in doc if page.get_text("text").strip()]
    doc.close()
    return "\n\n".join(pages), meta


def _read_docx(path: Path) -> tuple[str, dict]:
    from docx import Document
    doc = Document(str(path))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                paragraphs.append(row_text)
    return "\n\n".join(paragraphs), {}


def _read_pptx(path: Path) -> tuple[str, dict]:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [s.text.strip() for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides), {}


def _read_excel(path: Path) -> tuple[str, dict]:
    import pandas as pd
    ext = path.suffix.lower()
    df_dict = (
        {"Sheet1": pd.read_csv(str(path))} if ext == ".csv"
        else pd.read_excel(str(path), sheet_name=None, engine="openpyxl")
    )
    sheets = [f"[Sheet: {name}]\n{df.to_string(index=False)}" for name, df in df_dict.items()]
    return "\n\n".join(sheets), {}


def _read_text(path: Path) -> tuple[str, dict]:
    for enc in ("utf-8", "cp949", "euc-kr", "latin-1"):
        try:
            return path.read_text(encoding=enc), {}
        except UnicodeDecodeError:
            continue
    return path.read_bytes().decode("utf-8", errors="replace"), {}


def _read_image(path: Path, api_key: Optional[str] = None) -> tuple[str, dict]:
    """Claude Vision으로 이미지 OCR + 해석.
    모델은 src.config.models.get_vision_model()에서 자동 선택.
    """
    import anthropic

    ext_map = {
        ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
        ".png": "image/png",  ".webp": "image/webp",
        ".gif": "image/gif",  ".bmp": "image/png",
    }
    media_type = ext_map.get(path.suffix.lower(), "image/png")
    image_data = base64.standard_b64encode(path.read_bytes()).decode("utf-8")

    _, vision_model = get_vision_model()
    _log.debug(f"이미지 OCR 모델: {vision_model}")

    # 이미지 OCR은 Claude Vision 전용 (텍스트 LLM 폴백과 별개). 크레딧/키 문제로
    # 실패해도 업로드 흐름이 죽지 않도록 graceful 처리 — 빈 텍스트 + 경고 반환. (규칙10)
    try:
        client = anthropic.Anthropic(api_key=api_key or os.environ.get("ANTHROPIC_API_KEY"))
        response = client.messages.create(
            model=vision_model,
            max_tokens=2048,
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": media_type,
                            "data": image_data,
                        },
                    },
                    {
                        "type": "text",
                        "text": (
                            "Please extract ALL text visible in this image (OCR). "
                            "Then provide a concise interpretation of what the image shows "
                            "(charts, diagrams, tables, figures, etc.). "
                            "Format: first the extracted text, then '--- Interpretation ---', "
                            "then your interpretation."
                        ),
                    },
                ],
            }],
        )
        text = response.content[0].text if response.content else ""
        return text, {"vision_processed": True, "vision_model": vision_model}
    except Exception as e:
        _log.warning(
            "이미지 OCR 실패 (Claude Vision 전용 — 크레딧/키 필요). "
            "텍스트 논문/PDF/붙여넣기는 무료 작동: %s", str(e)[:120],
        )
        return "", {"vision_processed": False, "ocr_error": str(e)[:120]}


# ── Constants ─────────────────────────────────────────────────────────────────

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".pptx",
    ".xlsx", ".xls", ".csv",
    ".txt", ".md", ".rst",
    ".jpg", ".jpeg", ".png", ".webp", ".bmp", ".gif",
}

IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".gif"}


# ── Main class ────────────────────────────────────────────────────────────────

class DocumentReader:
    """모든 지원 형식을 읽어 통일된 document dict 반환."""

    def __init__(self, api_key: Optional[str] = None):
        self._api_key = api_key or os.environ.get("ANTHROPIC_API_KEY")

    def read(self, file_path: str) -> Dict:
        """파일을 읽어 document dict 반환.

        Returns:
            {path, filename, title, full_text, page_count, file_type, metadata}
        """
        path = Path(file_path).resolve()
        if not path.exists():
            raise FileNotFoundError(f"파일을 찾을 수 없습니다: {path}")

        ext = path.suffix.lower()
        if ext not in SUPPORTED_EXTENSIONS:
            raise ValueError(
                f"지원하지 않는 확장자 '{ext}'.\n"
                f"지원 형식: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
            )

        if ext == ".pdf":
            full_text, meta = _read_pdf(path)
            file_type = "pdf"
        elif ext == ".docx":
            full_text, meta = _read_docx(path)
            file_type = "word"
        elif ext == ".pptx":
            full_text, meta = _read_pptx(path)
            file_type = "powerpoint"
        elif ext in {".xlsx", ".xls", ".csv"}:
            full_text, meta = _read_excel(path)
            file_type = "spreadsheet"
        elif ext in {".txt", ".md", ".rst"}:
            full_text, meta = _read_text(path)
            file_type = "text"
        elif ext in IMAGE_EXTENSIONS:
            full_text, meta = _read_image(path, self._api_key)
            file_type = "image"
        else:
            raise ValueError(f"처리되지 않은 확장자: {ext}")

        title = meta.get("title") or path.stem
        word_count = len(full_text.split())
        page_count = max(1, word_count // 400)

        return {
            "path": str(path),
            "filename": path.name,
            "title": title,
            "full_text": full_text,
            "page_count": page_count,
            "file_type": file_type,
            "metadata": meta,
        }

    def read_directory(self, directory: str) -> List[Dict]:
        """디렉토리의 모든 지원 파일 읽기 (비재귀)."""
        dir_path = Path(directory)
        if not dir_path.is_dir():
            raise NotADirectoryError(f"디렉토리가 아닙니다: {dir_path}")
        results = []
        for f in sorted(dir_path.iterdir()):
            if f.suffix.lower() in SUPPORTED_EXTENSIONS:
                try:
                    results.append(self.read(str(f)))
                except Exception as exc:
                    _log.warning(f"[DocumentReader] 스킵: {f.name} — {exc}")
        return results

    def read_directory_recursive(self, directory: str) -> List[Dict]:
        """디렉토리의 모든 지원 파일 재귀 읽기."""
        dir_path = Path(directory)
        results = []
        for f in sorted(dir_path.rglob("*")):
            if f.is_file() and f.suffix.lower() in SUPPORTED_EXTENSIONS:
                try:
                    results.append(self.read(str(f)))
                except Exception as exc:
                    _log.warning(f"[DocumentReader] 스킵: {f.name} — {exc}")
        return results
